import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml import parse_xml

# ── Styling Constants ─────────────────────────────────────────────────────────
DARK_GREEN = RGBColor(45, 106, 79)    # #2D6A4F
OLIVE = RGBColor(75, 83, 32)         # #4B5320
GOLD = RGBColor(212, 160, 23)        # #D4A017
CHARCOAL = RGBColor(33, 37, 41)       # #212529
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 246, 248)  # #F5F6F8
MUTED_GRAY = RGBColor(108, 117, 125)  # #6C757D
PASS_GREEN = RGBColor(46, 125, 50)    # #2E7D32
FAIL_RED = RGBColor(198, 40, 40)      # #C62828

# ── Widescreen Aspect Ratio ───────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def init_widescreen_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


# ── Animation Transitions Helper ──────────────────────────────────────────────

def set_slide_transition_fade(slide):
    """
    Applies a smooth fade transition animation to the slide using OpenXML.
    """
    try:
        transition_xml = (
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:fade/>'
            '</p:transition>'
        )
        transition_el = parse_xml(transition_xml)
        slide.element.append(transition_el)
    except Exception:
        pass


# ── Presentation Layout Helpers ───────────────────────────────────────────────

def add_title_slide(prs, main_title, subtitle=""):
    """
    Creates a premium Army-styled widescreen Title Slide with animation transition.
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_transition_fade(slide)
    
    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_GREEN
    
    # Top Gold accent line
    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    # Title & Subtitle box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = main_title.upper()
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = WHITE
        p2.font.name = "Arial"
        p2.space_before = Pt(20)
        p2.alignment = PP_ALIGN.LEFT

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.8))
    tf_f = footer_box.text_frame
    tf_f.margin_left = Inches(0)
    tf_f.margin_right = Inches(0)
    p_f = tf_f.paragraphs[0]
    p_f.text = f"CONFIDENTIAL  |  INDIAN ARMY EVALUATION PORTAL  |  GENERATED ON {datetime.now().strftime('%d %b %Y').upper()}"
    p_f.font.size = Pt(11)
    p_f.font.bold = True
    p_f.font.color.rgb = RGBColor(160, 180, 160)
    p_f.font.name = "Arial"
    
    return slide


def add_base_slide(prs, title):
    """
    Creates a standard content slide with top header bar, gold divider, and footer.
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_transition_fade(slide)
    
    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Top header bar
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.8)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_GREEN
    top_bar.line.fill.background()
    
    # Gold divider line below top bar
    div_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), Inches(13.333), Inches(0.04)
    )
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = GOLD
    div_line.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12.0), Inches(0.6))
    tf = title_box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Arial"
    
    # Footer line & text
    footer_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4)
    )
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = RGBColor(240, 240, 240)
    footer_bg.line.fill.background()
    
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.12), Inches(12.0), Inches(0.35))
    tf_f = footer_box.text_frame
    tf_f.margin_left = Inches(0)
    tf_f.margin_top = Inches(0)
    p_f = tf_f.paragraphs[0]
    p_f.text = "CONFIDENTIAL - FOR OFFICIAL USE ONLY"
    p_f.font.size = Pt(9)
    p_f.font.bold = True
    p_f.font.color.rgb = MUTED_GRAY
    p_f.font.name = "Arial"
    
    return slide


def add_bullet_explanations(slide, x, y, cx, cy, points):
    """
    Creates an explanation text box on the right card with proper padding and margins.
    """
    # Background card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(215, 220, 225)
    card.line.width = Pt(1)
    
    # Text box (using generous inset margins to avoid any clipping/overlaps)
    tx_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), cx - Inches(0.6), cy - Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p_head = tf.paragraphs[0]
    p_head.text = "ANALYSIS & DIAGNOSTIC INSIGHTS"
    p_head.font.size = Pt(13)
    p_head.font.bold = True
    p_head.font.color.rgb = DARK_GREEN
    p_head.font.name = "Arial"
    p_head.space_after = Pt(12)
    
    for pt in points:
        p = tf.add_paragraph()
        p.text = f"•  {pt}"
        p.font.size = Pt(11)
        p.font.color.rgb = CHARCOAL
        p.font.name = "Arial"
        # Spacing to separate bullet lines cleanly
        p.space_after = Pt(10)
        
    return card


# ── Chart Slide Builders ──────────────────────────────────────────────────────

def add_donut_chart_slide(prs, title, categories, values, points):
    """
    Slide with a Doughnut chart on the left, explanations on the right.
    """
    slide = add_base_slide(prs, title)
    
    # Chart Data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Overall Results", tuple(values))
    
    # Left Position (Leaves proper margins around elements)
    x, y, cx, cy = Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.3)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    chart.legend.include_in_layout = False
    
    # Enable Data Labels to show values clearly on slices
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(11)
        data_labels.font.bold = True
        data_labels.font.color.rgb = WHITE
    except Exception:
        pass

    # Style segments (Green for Pass, Red for Fail)
    try:
        points_list = chart.series[0].points
        if len(points_list) >= 2:
            # Assumes index 0 = Passed, index 1 = Failed
            points_list[0].format.fill.solid()
            points_list[0].format.fill.fore_color.rgb = PASS_GREEN
            points_list[1].format.fill.solid()
            points_list[1].format.fill.fore_color.rgb = FAIL_RED
    except Exception:
        pass
        
    # Right Position Explanations
    add_bullet_explanations(
        slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.3), points
    )
    return slide


def add_bar_chart_slide(prs, title, categories, series_dict, points, is_stacked=False):
    """
    Slide with a native Column/Bar chart on the left, explanations on the right.
    """
    slide = add_base_slide(prs, title)
    
    # Chart Data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, vals in series_dict.items():
        chart_data.add_series(series_name, tuple(vals))
        
    chart_type = XL_CHART_TYPE.COLUMN_STACKED if is_stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
    
    x, y, cx, cy = Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.3)
    chart_shape = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9.5)
    
    # Enable Data Labels to show values clearly on columns
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.bold = True
        data_labels.font.color.rgb = CHARCOAL
    except Exception:
        pass

    # Try custom coloring series (Passed = Green, Failed = Red, Avg = Gold)
    try:
        for idx, series in enumerate(chart.series):
            if "pass" in series.name.lower() or "passed" in series.name.lower():
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = PASS_GREEN
            elif "fail" in series.name.lower() or "failed" in series.name.lower():
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = FAIL_RED
            elif "avg" in series.name.lower() or "marks" in series.name.lower() or "rate" in series.name.lower():
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = GOLD
    except Exception:
        pass
        
    # Right Position Explanations
    add_bullet_explanations(
        slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.3), points
    )
    return slide


def add_line_chart_slide(prs, title, categories, series_dict, points):
    """
    Slide with a native Line chart on the left, explanations on the right.
    """
    slide = add_base_slide(prs, title)
    
    # Chart Data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, vals in series_dict.items():
        chart_data.add_series(series_name, tuple(vals))
        
    x, y, cx, cy = Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.3)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9.5)
    
    # Enable Data Labels to show monthly trends values
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True
        data_labels.font.color.rgb = CHARCOAL
    except Exception:
        pass

    # Color lines
    try:
        for series in chart.series:
            if "count" in series.name.lower() or "total" in series.name.lower():
                series.format.line.color.rgb = DARK_GREEN
            elif "pass" in series.name.lower():
                series.format.line.color.rgb = PASS_GREEN
    except Exception:
        pass
        
    # Right Position Explanations
    add_bullet_explanations(
        slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.3), points
    )
    return slide


def add_table_slide(prs, title, headers, rows, points):
    """
    Slide with a beautifully styled Table on the left, explanations on the right.
    """
    slide = add_base_slide(prs, title)
    
    # Table bounds
    tx, ty, tcx, tcy = Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.3)
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, tx, ty, tcx, tcy)
    table = table_shape.table
    
    # Set uniform column widths
    col_width = int(tcx / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_width
    
    # Format Headers
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(text).upper()
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_GREEN
        
        # Adjust cell margins to avoid overlapping text
        tf = cell.text_frame
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = GOLD
            
    # Format Data Rows
    for row_idx, rdata in enumerate(rows):
        for col_idx, val in enumerate(rdata):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            # Alternate row background colors
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.fore_color.rgb = WHITE
                
            # Adjust cell margins
            tf = cell.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(9.5)
                p.font.color.rgb = CHARCOAL
                
    # Right Position Explanations
    add_bullet_explanations(
        slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.3), points
    )
    return slide


# ── Presentation Builders (Export Entry Points) ───────────────────────────────

def generate_commander_pptx(data):
    """
    Generates presentation object for Commander Dashboard.
    """
    prs = init_widescreen_presentation()
    add_title_slide(
        prs, 
        "Commander's Dashboard Review", 
        "Comprehensive Evaluation & Operational Readiness Report"
    )
    
    # Slide 2: Overall metrics summary
    stats_rows = [
        ["Metric", "Count / Value"],
        ["Total Agniveers", data['total_agniveers']],
        ["Evaluated Agniveers", data['evaluated_agniveers']],
        ["Qualified (Pass)", data['pass_count']],
        ["Unqualified (Fail)", data['fail_count']],
        ["Passing Percentage", f"{data['pass_rate']}%"],
        ["Course Completion Rate", f"{data['completion_rate']}%"],
    ]
    summary_bullets = [
        f"A total of {data['total_agniveers']} Agniveers are registered in the portal database.",
        f"Currently, {data['evaluated_agniveers']} trainees have completed their full evaluation cycle.",
        f"A total of {data['pass_count']} trainees qualified the standards, yielding a passing rate of {data['pass_rate']}%.",
        f"{data['fail_count']} trainees did not qualify the required standards and are flagged for corrective coaching.",
        f"The training batch has achieved an overall completion rate of {data['completion_rate']}% relative to strength."
    ]
    add_table_slide(
        prs, 
        "Executive Summary - Batch Metrics", 
        stats_rows[0], 
        stats_rows[1:], 
        summary_bullets
    )
    
    # Slide 3: Pass/Fail Ratio
    ratio_bullets = [
        f"Out of {data['evaluated_agniveers']} evaluated Agniveers, {data['pass_count']} qualified (Pass) and {data['fail_count']} failed to qualify (Fail).",
        f"The overall passing rate is {data['pass_rate']}% while the failure rate stands at {100 - float(data['pass_rate']):.1f}%.",
        "Passing is determined by achieving a cumulative percentage >= 50% across departments (40% in Battalion specific tests).",
        "Target is to reach a 100% operational qualification level. Revision training is structured for unqualified personnel."
    ]
    add_donut_chart_slide(
        prs,
        "Overall Passing Status Ratio",
        ["Passed", "Failed"],
        [data['pass_count'], data['fail_count']],
        ratio_bullets
    )
    
    # Slide 4: Department-wise Breakdown
    dept_labels = data['dept_labels']
    dept_passed = data['dept_passed']
    dept_failed = data['dept_failed']
    series_dict = {
        "Passed": dept_passed,
        "Failed": dept_failed
    }
    
    dept_bullets = [
        "Comparative trainee evaluation analysis across training wings:",
    ]
    for i, name in enumerate(dept_labels):
        p = dept_passed[i]
        f = dept_failed[i]
        t = p + f
        rate = (p / max(t, 1)) * 100
        dept_bullets.append(f"{name} Wing: {p} qualified, {f} unqualified out of {t} evaluated ({rate:.1f}% Pass rate).")
    dept_bullets.append("Battalion (A) covers physical & tactical; TTS (B) covers tech trades; CS (C) covers fieldcraft; Clerk (D) handles administration.")
    
    add_bar_chart_slide(
        prs,
        "Trainee Performance by Department",
        dept_labels,
        series_dict,
        dept_bullets
    )
    
    # Slide 5: Monthly Trend
    months = data['months_labels']
    counts = data['month_counts']
    pass_counts = data['month_pass']
    trend_dict = {
        "Total Evaluated": counts,
        "Passed Candidates": pass_counts
    }
    
    trend_bullets = [
        "Month-over-month training progress and qualification trends:",
    ]
    for i, m in enumerate(months):
        c = counts[i]
        p = pass_counts[i]
        rate = (p / max(c, 1)) * 100
        trend_bullets.append(f"{m}: {c} evaluated, {p} qualified ({rate:.1f}% monthly pass rate).")
    trend_bullets.append("The monthly volume variance corresponds to batch intakes and course graduation timelines.")
    
    add_line_chart_slide(
        prs,
        "Monthly Progress & Evaluation Trends",
        months,
        trend_dict,
        trend_bullets
    )
    
    # Slide 6: Top Performers Table
    top_performers = data.get('top_agniveers', [])
    top_rows = []
    for ag in top_performers[:8]:
        top_rows.append([
            ag.get('enrollment', '—'),
            ag.get('name', '—'),
            f"{ag.get('percentage', 0)}%"
        ])
    top_bullets = [
        "Lists the outstanding trainees who secured the highest overall scores.",
        "These individuals display exceptional proficiency across all military subjects.",
        "Merit listings are used to select batch commanders and drill leaders.",
        "Recommended for early leadership roles, key assignments, and recognition."
    ]
    add_table_slide(
        prs,
        "Top Performers Merit List",
        ["Enrollment No", "Name", "Score (%)"],
        top_rows if top_rows else [["—", "No data", "—"]],
        top_bullets
    )
    
    return prs


def generate_ghead_pptx(data):
    """
    Generates presentation object for G-Head Dashboard.
    """
    prs = init_widescreen_presentation()
    add_title_slide(
        prs, 
        "G-Head Dashboard Review", 
        "Training Quality Assurance & Subject Performance Review"
    )
    
    # Slide 2: Overall Metrics
    stats_rows = [
        ["Metric", "Count / Value"],
        ["Total Agniveers", data['total_agniveers']],
        ["Evaluated Agniveers", data['evaluated_agniveers']],
        ["Passed", data['pass_count']],
        ["Failed", data['fail_count']],
        ["Passing Percentage", f"{data['pass_rate']}%"],
    ]
    summary_bullets = [
        f"Total trainee strength registered in database is {data['total_agniveers']}.",
        f"Evaluations completed for {data['evaluated_agniveers']} trainees.",
        f"The training cycle shows a positive passing rate of {data['pass_rate']}% across modules.",
        f"Currently, {data['fail_count']} trainees have not met the passing standard."
    ]
    add_table_slide(
        prs, 
        "Training Executive Summary", 
        stats_rows[0], 
        stats_rows[1:], 
        summary_bullets
    )

    # Slide 3: Overall Ratio
    ratio_bullets = [
        f"Pass/Fail distribution among evaluated trainees: {data['pass_rate']}% qualified vs {100 - float(data['pass_rate']):.1f}% failed.",
        "Represents the overall qualification status across all training modules.",
        "Corrective training is recommended for the failed cohort before final deployment."
    ]
    add_donut_chart_slide(
        prs,
        "Overall Qualification Status",
        ["Passed", "Failed"],
        [data['pass_count'], data['fail_count']],
        ratio_bullets
    )

    # Slide 4: Category Pass Rates (Physical, Drill, Weapon, Written, etc.)
    cat_labels = data['category_labels']
    cat_rates = data['category_pass_rates']
    
    cat_bullets = [
        "Passing percentage across core subject domains:",
    ]
    for i, label in enumerate(cat_labels):
        rate = cat_rates[i]
        cat_bullets.append(f"{label}: {rate}% pass rate among evaluated candidates.")
    cat_bullets.append("Skill domains with pass rates below 75% are prioritized for instructor review and extra practice hours.")
    
    add_bar_chart_slide(
        prs,
        "Pass Rates by Subject Category",
        cat_labels,
        {"Pass Rate (%)": cat_rates},
        cat_bullets
    )

    # Slide 5: Department Breakdown
    dept_labels = data['dept_labels']
    dept_passed = data['dept_passed']
    dept_failed = data['dept_failed']
    series_dict = {
        "Passed": dept_passed,
        "Failed": dept_failed
    }
    
    dept_bullets = [
        "Comparative analysis across specialized training wings:",
    ]
    for i, name in enumerate(dept_labels):
        p = dept_passed[i]
        f = dept_failed[i]
        t = p + f
        rate = (p / max(t, 1)) * 100
        dept_bullets.append(f"{name} Wing: {p} passed, {f} failed out of {t} evaluated ({rate:.1f}% Pass rate).")
    
    add_bar_chart_slide(
        prs,
        "Trainee Status by Department",
        dept_labels,
        series_dict,
        dept_bullets
    )
    
    # Slide 6: Monthly Trend
    months = data['months_labels']
    counts = data['month_counts']
    pass_counts = data['month_pass']
    trend_dict = {
        "Total Evaluated": counts,
        "Passed Candidates": pass_counts
    }
    
    trend_bullets = [
        "Month-over-month training progression and evaluation trends:",
    ]
    for i, m in enumerate(months):
        c = counts[i]
        p = pass_counts[i]
        rate = (p / max(c, 1)) * 100
        trend_bullets.append(f"{m}: {c} evaluated, {p} qualified ({rate:.1f}% Pass rate).")
    
    add_line_chart_slide(
        prs,
        "Monthly Training Progression Trends",
        months,
        trend_dict,
        trend_bullets
    )
    
    return prs


def generate_department_pptx(dept_code, data):
    """
    Generates presentation object for Department-specific Dashboard.
    """
    prs = init_widescreen_presentation()
    dept_names = {'A': 'Battalion', 'B': 'TTS', 'C': 'CS', 'D': 'Clerk'}
    dept_name = dept_names.get(dept_code, f"Department {dept_code}")
    
    add_title_slide(
        prs, 
        f"Department {dept_name.upper()} Review", 
        f"Trainee Subject Performance & Skill Qualification Report"
    )
    
    # Slide 2: Department Summary Table
    stats_rows = [
        ["Metric", "Count / Value"],
        ["Total Agniveers In Scope", data['total_agniveers']],
        ["Evaluated Trainees", data['evaluated_agniveers']],
        ["Passed Standard", data['pass_count']],
        ["Failed Standard", data['fail_count']],
        ["Passing Rate", f"{data['pass_rate']}%"],
        ["Completion Rate", f"{data['completion_rate']}%"],
    ]
    summary_bullets = [
        f"This slide summaries the core metrics for Department {dept_name}.",
        f"Out of {data['total_agniveers']} assigned trainees, {data['evaluated_agniveers']} have completed their evaluation sheets.",
        f"The passing rate inside this department stands at {data['pass_rate']}%.",
        f"Trainee syllabus completion rate currently stands at {data['completion_rate']}%."
    ]
    add_table_slide(
        prs, 
        f"{dept_name} Executive Summary", 
        stats_rows[0], 
        stats_rows[1:], 
        summary_bullets
    )

    # Slide 3: Pass/Fail Donut
    ratio_bullets = [
        f"Shows the qualification ratio in {dept_name}: {data['pass_rate']}% passed, {100-float(data['pass_rate']):.1f}% failed.",
        "Trainees must qualify all required parameters to pass the overall department standard.",
        "Target is to raise the qualification rate through supervised practice and re-evaluation sessions."
    ]
    add_donut_chart_slide(
        prs,
        "Qualification Status Distribution",
        ["Passed", "Failed"],
        [data['pass_count'], data['fail_count']],
        ratio_bullets
    )

    # Slide 4: Subject Category Pass Rates
    cat_labels = data['category_bar_labels']
    cat_rates = data['category_pass_rates']
    
    cat_bullets = [
        "Department skill category qualification rates:",
    ]
    for i, label in enumerate(cat_labels):
        rate = cat_rates[i]
        cat_bullets.append(f"{label}: {rate}% pass rate achieved in this category.")
    cat_bullets.append("Instructions are adjusted for categories showing pass rates lower than the 75% target.")
    
    add_bar_chart_slide(
        prs,
        "Skill Category Qualification Rates",
        cat_labels,
        {"Pass Rate (%)": cat_rates},
        cat_bullets
    )

    # Slide 5: Test Type Performance
    test_labels = data['test_labels']
    avg_marks = data['test_avg_marks']
    pass_rates = data.get('test_pass_rates', [])
    
    test_bullets = [
        "Detailed test-wise evaluation statistics:",
    ]
    for i, label in enumerate(test_labels):
        avg = avg_marks[i]
        pr = pass_rates[i] if i < len(pass_rates) else 0.0
        test_bullets.append(f"{label}: Avg score of {avg:g} marks with a {pr:.1f}% pass rate.")
    
    add_bar_chart_slide(
        prs,
        "Test-Wise Average Scores",
        test_labels,
        {"Average Score": avg_marks},
        test_bullets
    )

    # Slide 6: Monthly Trend
    months = data['months_labels']
    counts = data['month_counts']
    pass_counts = data['month_pass']
    
    trend_bullets = [
        "Monthly department progress and evaluation trends:",
    ]
    for i, m in enumerate(months):
        c = counts[i]
        p = pass_counts[i]
        rate = (p / max(c, 1)) * 100
        trend_bullets.append(f"{m}: {c} evaluated, {p} qualified ({rate:.1f}% Pass rate).")
    
    add_line_chart_slide(
        prs,
        "Monthly Department Progress Trend",
        months,
        {"Evaluated": counts, "Passed": pass_counts},
        trend_bullets
    )

    return prs


def generate_registration_pptx(data):
    """
    Generates presentation object for Registration Dashboard.
    """
    prs = init_widescreen_presentation()
    add_title_slide(
        prs, 
        "Agniveer Registration Summary", 
        "Enrollment, Batch Sizes, Trade Distribution & Documentation Report"
    )
    
    # Slide 2: Overall Metrics
    stats_rows = [
        ["Parameter", "Details"],
        ["Total Agniveers Registered", data['total_count']],
        ["Active Batch", data.get('active_batch', 'Batch 2026')],
        ["Date Compiled", datetime.now().strftime('%d %b %Y')],
    ]
    summary_bullets = [
        f"Overall total of registered Agniveers in the portal is {data['total_count']}.",
        "Trainees are enrolled individually or imported in batches via Excel sheets.",
        "The registration office holds primary responsibility for verifying all entry documents."
    ]
    add_table_slide(
        prs, 
        "Registration Overview", 
        stats_rows[0], 
        stats_rows[1:], 
        summary_bullets
    )

    # Slide 3: Trade-wise Distribution
    trades_dict = data.get('trades_dict', {})
    trade_labels = list(trades_dict.keys())
    trade_values = list(trades_dict.values())
    total_count = data['total_count']
    
    trade_bullets = [
        "Trainee distribution across specialized trade streams:",
    ]
    for label, val in trades_dict.items():
        pct = (val / max(total_count, 1)) * 100
        trade_bullets.append(f"{label}: {val} trainees enrolled ({pct:.1f}% of total batch).")
    trade_bullets.append("Trades include General Duty (GD), Technical (Tech), Clerk (CLK), Motor Vehicle (DMV), and others.")
    
    if trade_labels:
        add_bar_chart_slide(
            prs,
            "Trainees Breakdown by Trade",
            trade_labels,
            {"Trainees Count": trade_values},
            trade_bullets
        )
    else:
        # Fallback table
        add_table_slide(
            prs,
            "Trainees Breakdown by Trade",
            ["Trade", "Count"],
            [["GD", "N/A"], ["Tech", "N/A"], ["Clerk", "N/A"]],
            trade_bullets
        )

    # Slide 4: Company / Platoon Distribution
    company_dict = data.get('company_dict', {})
    comp_labels = list(company_dict.keys())
    comp_values = list(company_dict.values())
    
    comp_bullets = [
        "Trainee distribution across active training companies:",
    ]
    for label, val in company_dict.items():
        pct = (val / max(total_count, 1)) * 100
        comp_bullets.append(f"{label}: {val} trainees assigned ({pct:.1f}% of total batch).")
    comp_bullets.append("Workload is distributed evenly among companies for drill and weapon training.")
    
    if comp_labels:
        add_bar_chart_slide(
            prs,
            "Trainees Breakdown by Company",
            comp_labels,
            {"Count": comp_values},
            comp_bullets
        )
    else:
        add_table_slide(
            prs,
            "Trainees Breakdown by Company",
            ["Company", "Count"],
            [["A Coy", "N/A"], ["B Coy", "N/A"], ["C Coy", "N/A"], ["D Coy", "N/A"]],
            comp_bullets
        )

    # Slide 5: Recent Registrations Activity
    recent_list = data.get('recent_registrations', [])
    recent_rows = []
    for r in recent_list[:8]:
        recent_rows.append([
            r.agniveer_no or r.enrollment_number or '—',
            r.name,
            r.trade or '—',
            r.created_at.strftime('%Y-%m-%d') if r.created_at else '—'
        ])
        
    recent_bullets = [
        "Logs the latest trainee records added to the database.",
        "Useful for auditing daily registration activities.",
        "Allows registration officers to verify that correct enrollment numbers are generated."
    ]
    add_table_slide(
        prs,
        "Recent Registration Activity Log",
        ["Enrollment No", "Name", "Trade", "Date Registered"],
        recent_rows if recent_rows else [["—", "No data", "—", "—"]],
        recent_bullets
    )
    
    return prs
